"""Parser for PowerPoint (.pptx) files using python-pptx.

Extracts text from text frames and tables page-by-page (slide-by-slide).
Also extracts images from slide shapes and saves them to the wiki/assets folder.
"""

from __future__ import annotations

import os
from pathlib import Path
from .base import (
    DocumentParser,
    ParsedDocument,
    ParserError,
    compute_hash,
    fallback_title_from_path,
    normalize_text,
)


class PowerpointParser(DocumentParser):
    def can_parse(self, suffix: str) -> bool:
        return suffix.lower() == ".pptx"

    def parse(self, path: Path, chunk_size: int = 1500, overlap: int = 200) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ParserError(
                "python-pptx is not installed. Run `uv pip install -e .` to install dependencies."
            ) from e

        try:
            prs = Presentation(str(path))
        except Exception as e:
            raise ParserError(f"Cannot open PPTX {path.name}: {e}") from e

        # Resolve wiki assets folder
        from .. import config as cfg
        root = cfg.find_wiki_root(path) or path.parent
        paths = cfg.WikiPaths(root=root)
        assets_dir = paths.wiki / "assets"
        assets_dir.mkdir(parents=True, exist_ok=True)

        slide_chunks: list[str] = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_text_lines: list[str] = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        txt = paragraph.text.strip()
                        if txt:
                            slide_text_lines.append(txt)
                
                # Extract table text if present
                if shape.has_table:
                    for row in shape.table.rows:
                        row_txt = [cell.text.strip() for cell in row.cells]
                        if any(row_txt):
                            slide_text_lines.append(" | ".join(row_txt))

            slide_text = "\n".join(slide_text_lines)
            
            # Extract pictures/images
            images_markdown = []
            image_idx = 1
            for shape in slide.shapes:
                if hasattr(shape, "image") and shape.image:
                    try:
                        image_file_object = shape.image
                        ext = ".png"
                        if image_file_object.content_type:
                            if "jpeg" in image_file_object.content_type:
                                ext = ".jpg"
                            elif "png" in image_file_object.content_type:
                                ext = ".png"
                            elif "gif" in image_file_object.content_type:
                                ext = ".gif"
                        
                        img_name = f"{path.stem}_slide_{slide_idx + 1}_{image_idx}{ext}"
                        img_path = assets_dir / img_name
                        with open(img_path, "wb") as f:
                            f.write(image_file_object.blob)
                        
                        images_markdown.append(f"![Extracted Slide {slide_idx + 1} Image {image_idx}](assets/{img_name})")
                        image_idx += 1
                    except Exception:
                        pass

            if images_markdown:
                slide_text += "\n\n" + "\n".join(images_markdown)

            if slide_text.strip() or images_markdown:
                slide_chunks.append(normalize_text(slide_text))

        full_text = "\n\n".join(slide_chunks)
        text = normalize_text(full_text)
        title = fallback_title_from_path(path)
        metadata = {"slide_count": len(prs.slides)}

        return ParsedDocument(
            source_path=path,
            file_type="pptx",
            title=title,
            text=text,
            content_hash=compute_hash(text),
            bytes=path.stat().st_size,
            metadata=metadata,
            chunks=slide_chunks,
        )


def parse(path: Path) -> ParsedDocument:
    """Parse a .pptx file."""
    return PowerpointParser().parse(path)
