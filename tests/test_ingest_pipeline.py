"""Phase 1 ingest pipeline and unsupported-input guard tests."""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from llm_wiki.cli import build_parser
from llm_wiki.pipeline import (
    UnsupportedInputError,
    UserInputError,
    ingest_markdown_file,
)
from llm_wiki.pipeline.hashing import (
    MARKDOWN_SUFFIXES,
    UNSUPPORTED_SUFFIX_GUIDANCE,
    validate_markdown_input,
)
from llm_wiki.workspace import resolve_workspace


def _invoke(cli_args: list[str], path: Path) -> tuple[int, dict[str, object]]:
    argv = [*cli_args, "--path", str(path), "--json"]
    parser = build_parser()
    args = parser.parse_args(argv)
    return args.handler(args)


def _ensure_init(workspace: Path) -> None:
    _invoke(["init"], workspace)


def test_validate_markdown_input_accepts_supported_suffix() -> None:
    validate_markdown_input("/tmp/anything.md")
    validate_markdown_input("/tmp/anything.markdown")


def test_validate_markdown_input_rejects_unsupported_suffix() -> None:
    for suffix in UNSUPPORTED_SUFFIX_GUIDANCE:
        with pytest.raises(UnsupportedInputError):
            validate_markdown_input(f"/tmp/sample{suffix}")


def test_validate_markdown_input_rejects_url() -> None:
    with pytest.raises(UnsupportedInputError):
        validate_markdown_input("https://example.com/article")
    with pytest.raises(UnsupportedInputError):
        validate_markdown_input("http://example.com/")


def test_validate_markdown_input_rejects_unknown_suffix() -> None:
    with pytest.raises(UserInputError):
        validate_markdown_input("/tmp/sample.txt")


def test_ingest_creates_source_row_and_wiki_outputs(workspace: Path, samples_dir: Path) -> None:
    _ensure_init(workspace)
    sample = samples_dir / "short-note.md"
    assert sample.exists()

    exit_code, payload = _invoke(["ingest", str(sample)], workspace)
    assert exit_code == 0
    assert payload["status"] == "ok"
    source_id = payload["source_id"]
    assert source_id.startswith("source_")

    # Source summary compatibility file lives under vault/10_Wiki/sources/.
    stub_path = workspace / payload["source_stub_path"]
    assert stub_path.exists()
    stub_body = stub_path.read_text(encoding="utf-8")
    assert "Generated wiki pages" in stub_body
    assert payload["page_count"] >= 1
    assert payload["wiki_pages"]

    # After HI-3 fix, ingest no longer pre-creates a queued normalize job
    # (`wiki normalize` owns the canonical normalize job). ingest only
    # returns the artifact id; explicit `wiki normalize` is required.
    assert payload["job_id"] is None
    assert payload["artifact_id"].startswith("artifact_")
    artifact_path = workspace / payload["artifact_path"]
    assert artifact_path.exists()
    stored = json.loads(artifact_path.read_text(encoding="utf-8"))
    assert stored["status"] == "ok"
    assert stored["source_id"] == source_id
    assert stored["normalize_job_id"] is None


def test_ingest_malformed_pdf_returns_exit_code_2(workspace: Path, tmp_path: Path) -> None:
    _ensure_init(workspace)
    fake_pdf = tmp_path / "paper.pdf"
    fake_pdf.write_bytes(b"%PDF-1.4 fake")
    from llm_wiki.cli import main as cli_main

    exit_code = cli_main(["ingest", str(fake_pdf), "--path", str(workspace)])
    # Malformed/scanned PDFs without extractable text still fail cleanly.
    assert exit_code == 2


def test_ingest_txt_converts_to_markdown_pipeline(workspace: Path, tmp_path: Path) -> None:
    _ensure_init(workspace)
    text_file = tmp_path / "plain-note.txt"
    text_file.write_text(
        "Plain Note Title\n\nThis plain text document has enough useful words to flow through ingest normalize chunk embed wiki generation.",
        encoding="utf-8",
    )

    exit_code, payload = _invoke(["ingest", str(text_file)], workspace)

    assert exit_code == 0
    assert payload["status"] == "ok"
    raw_path = workspace / payload["raw_path"]
    assert raw_path.suffix == ".md"
    assert "# Plain Note Title" in raw_path.read_text(encoding="utf-8")


def test_ingest_docx_converts_to_markdown_pipeline(workspace: Path, tmp_path: Path) -> None:
    _ensure_init(workspace)
    from docx import Document

    docx_file = tmp_path / "docx-note.docx"
    document = Document()
    document.add_heading("DOCX Ingest Title", level=1)
    document.add_paragraph("This DOCX document has enough useful words to flow through ingest normalize chunk embed wiki generation.")
    document.save(docx_file)

    exit_code, payload = _invoke(["ingest", str(docx_file)], workspace)

    assert exit_code == 0
    assert payload["status"] == "ok"
    raw_path = workspace / payload["raw_path"]
    raw_text = raw_path.read_text(encoding="utf-8")
    assert raw_path.suffix == ".md"
    assert "# DOCX Ingest Title" in raw_text
    from llm_wiki.db.schema import connect

    conn = connect(resolve_workspace(workspace).db)
    try:
        row = conn.execute("SELECT source_type, metadata_json FROM sources WHERE id = ?", (payload["source_id"],)).fetchone()
    finally:
        conn.close()
    assert row["source_type"] == "converted_markdown"
    assert json.loads(row["metadata_json"])["conversion_source"] == "docx"

    duplicate_exit, duplicate_payload = _invoke(["ingest", str(docx_file)], workspace)
    assert duplicate_exit == 0
    assert duplicate_payload["status"] == "duplicate"
    assert duplicate_payload["source_id"] == payload["source_id"]


def test_ingest_pptx_converts_to_markdown_pipeline(workspace: Path, tmp_path: Path) -> None:
    _ensure_init(workspace)
    from pptx import Presentation

    pptx_file = tmp_path / "deck-note.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "PPTX Ingest Title"
    slide.placeholders[1].text = "This PowerPoint slide has enough useful words to flow through ingest normalize chunk embed wiki generation."
    presentation.save(pptx_file)

    exit_code, payload = _invoke(["ingest", str(pptx_file)], workspace)

    assert exit_code == 0
    assert payload["status"] == "ok"
    raw_text = (workspace / payload["raw_path"]).read_text(encoding="utf-8")
    assert "PPTX Ingest Title" in raw_text


def test_ingest_xlsx_converts_to_markdown_pipeline(workspace: Path, tmp_path: Path) -> None:
    _ensure_init(workspace)
    from openpyxl import Workbook

    xlsx_file = tmp_path / "sheet-note.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Knowledge"
    sheet.append(["Title", "Body"])
    sheet.append(["XLSX Ingest Title", "This spreadsheet row has enough useful words to flow through ingest normalize chunk embed wiki generation."])
    workbook.save(xlsx_file)

    exit_code, payload = _invoke(["ingest", str(xlsx_file)], workspace)

    assert exit_code == 0
    assert payload["status"] == "ok"
    raw_text = (workspace / payload["raw_path"]).read_text(encoding="utf-8")
    assert "XLSX Ingest Title" in raw_text


def test_scan_inbox_root_archives_successes_and_writes_review_markdown(workspace: Path, tmp_path: Path) -> None:
    _ensure_init(workspace)
    from pypdf import PdfWriter
    from pypdf.generic import DictionaryObject, NameObject, StreamObject
    from llm_wiki.pipeline.ingest import scan_inbox
    from llm_wiki.workspace import resolve_workspace
    import sqlite3

    paths = resolve_workspace(workspace)
    inbox = tmp_path / "00. Inbox"
    inbox.mkdir(parents=True)
    (inbox / "_Review").mkdir()
    (inbox / "_Failed").mkdir()
    ignored_child = inbox / "Files"
    ignored_child.mkdir()
    (ignored_child / "ignored.md").write_text("# Ignored child folder\n", encoding="utf-8")
    md_file = inbox / "direct-note.md"
    md_file.write_text("# Direct Note\n\nRoot markdown should be archived after successful ingest.\n", encoding="utf-8")
    pdf_file = inbox / "move-policy.pdf"

    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject({NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type1"), NameObject("/BaseFont"): NameObject("/Helvetica")})
    font_ref = writer._add_object(font)
    page[NameObject("/Resources")] = DictionaryObject({NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})})
    stream = StreamObject()
    unique_text = (
        "Inbox PDF Move Policy Unique Text 20260723. This generated text PDF contains enough words "
        "for MarkItDown pdfplumber extraction and verifies that successful inbox documents "
        "are copied to Review while originals are archived outside the Obsidian vault."
    )
    stream._data = f"BT /F1 18 Tf 72 720 Td ({unique_text}) Tj ET".encode("latin-1")
    page[NameObject("/Contents")] = writer._add_object(stream)
    with pdf_file.open("wb") as f:
        writer.write(f)

    payload = scan_inbox(paths, [inbox])

    assert payload["status"] == "ok"
    assert payload["new_candidate_count"] == 2
    assert len(payload["postprocessed"]) == 2
    assert not pdf_file.exists()
    assert not md_file.exists()
    assert (ignored_child / "ignored.md").exists()

    from llm_wiki.common import utc_now
    expected_date = utc_now()[:10].replace("-", "")
    review_paths = {Path(post["markdown_to"]).name: Path(post["markdown_to"]) for post in payload["postprocessed"]}
    assert f"move-policy_{expected_date}.md" in review_paths
    assert f"direct-note_{expected_date}.md" in review_paths
    for post in payload["postprocessed"]:
        archived_original = Path(post["original_to"])
        review_copy = Path(post["markdown_to"])
        assert archived_original.exists()
        assert paths.vault not in archived_original.parents
        assert archived_original.is_relative_to(paths.data / "inbox_originals")
        assert review_copy.exists()
        assert review_copy.parent == inbox / "_Review"
    assert "Inbox PDF Move Policy Unique Text 20260723" in review_paths[f"move-policy_{expected_date}.md"].read_text(encoding="utf-8")

    conn = sqlite3.connect(paths.db)
    conn.row_factory = sqlite3.Row
    try:
        rows = conn.execute("SELECT origin, metadata_json FROM sources ORDER BY created_at").fetchall()
    finally:
        conn.close()
    assert len(rows) == 2
    for row in rows:
        metadata = json.loads(row["metadata_json"])
        assert Path(row["origin"]).exists()
        assert "archived_original_path" in metadata
        assert "review_markdown_path" in metadata


def test_ingest_unsupported_url_returns_exit_code_2(workspace: Path) -> None:
    _ensure_init(workspace)
    from llm_wiki.cli import main as cli_main

    exit_code = cli_main(["ingest", "https://example.com/article", "--path", str(workspace)])
    assert exit_code == 2


def test_ingest_text_creates_user_text_source(workspace: Path) -> None:
    _ensure_init(workspace)
    body = "User typed knowledge from clipboard."

    from llm_wiki.cli import main as cli_main

    exit_code = cli_main(
        ["ingest-text", "Clip note", "--text", body, "--path", str(workspace)]
    )
    assert exit_code == 0


def test_ingest_duplicate_returns_duplicate_status(workspace: Path, samples_dir: Path) -> None:
    _ensure_init(workspace)
    sample = samples_dir / "short-note.md"

    first_exit, first_payload = _invoke(["ingest", str(sample)], workspace)
    assert first_exit == 0
    assert first_payload["status"] == "ok"
    first_source_id = first_payload["source_id"]

    # ingest_markdown_file is the underlying pipeline call; pass the same path
    # in-process to confirm the duplicate code branch without spawning another
    # subprocess.
    paths = resolve_workspace(workspace)
    result = ingest_markdown_file(paths, sample)
    assert result["status"] == "duplicate"
    assert result["source_id"] == first_source_id
